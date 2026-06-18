import os
import re
from fastapi import HTTPException

def extract_text(file_path: str) -> str:
    """
    Extract raw text from documents. Supports PDF, DOCX, XLSX, PPTX, and standard text/code files.
    """
    _, ext = os.path.splitext(file_path.lower())
    
    # 1. Plain text and code files
    if ext in [".txt", ".md", ".json", ".csv", ".py", ".js", ".ts", ".tsx", ".html", ".css", ".yaml", ".yml", ".ini", ".conf"]:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to read text file: {str(e)}")
            
    # 2. PDF Files
    elif ext == ".pdf":
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            text_list = [page.extract_text() or "" for page in reader.pages]
            return "\n".join(text_list)
        except ImportError:
            raise HTTPException(status_code=500, detail="PDF parser 'pypdf' is not installed.")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse PDF: {str(e)}")
            
    # 3. Word Files (DOCX)
    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_content = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_content:
                        full_text.append(" | ".join(row_content))
            return "\n".join(full_text)
        except ImportError:
            raise HTTPException(status_code=500, detail="Word parser 'python-docx' is not installed.")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse Word Document: {str(e)}")
            
    # 4. Excel Files (XLSX)
    elif ext == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_list = []
            for sheet in wb.sheetnames:
                text_list.append(f"--- Sheet: {sheet} ---")
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(val).strip() if val is not None else "" for val in row])
                    # Clean up multiple empty fields
                    row_text = re.sub(r"(\s*\|\s*){2,}", " | ", row_text).strip(" | ")
                    if row_text:
                        text_list.append(row_text)
            return "\n".join(text_list)
        except ImportError:
            raise HTTPException(status_code=500, detail="Excel parser 'openpyxl' is not installed.")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse Excel file: {str(e)}")

    # 5. PowerPoint Files (PPTX)
    elif ext == ".pptx":
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            text_list = []
            for i, slide in enumerate(prs.slides):
                text_list.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_list.append(shape.text.strip())
            return "\n".join(text_list)
        except ImportError:
            raise HTTPException(status_code=500, detail="PowerPoint parser 'python-pptx' is not installed.")
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to parse PowerPoint presentation: {str(e)}")

    else:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file format '{ext}' for RAG parsing."
        )

def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> list:
    """
    Split text into overlapping chunks.
    """
    chunks = []
    # Normalize multiple newlines/whitespaces to single space for clean vector parsing
    normalized = re.sub(r"\s+", " ", text).strip()
    
    start = 0
    text_len = len(normalized)
    
    if text_len <= chunk_size:
        return [normalized] if normalized else []
        
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunk = normalized[start:end]
        if chunk.strip():
            chunks.append(chunk)
        start += (chunk_size - overlap)
        
    return chunks
