import os
import re
import json
import logging
from typing import Tuple

logger = logging.getLogger(__name__)

# Try optional imports
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

# Local storage path mapping
UPLOAD_DIR = os.getenv("VOXKAGE_UPLOAD_DIR", "/app/data/uploads")
if not os.path.exists(UPLOAD_DIR):
    os.makedirs(UPLOAD_DIR, exist_ok=True)

def _resolve_server_path(filename: str, directory: str = "") -> str:
    """
    Resolves a file path on the server. Map standard user dirs (Downloads, Desktop, etc.)
    directly to the server's uploads folder.
    """
    # Clean up filename
    filename = os.path.basename(filename)
    
    # If no directory or standard desktop directory, put it directly in UPLOAD_DIR
    clean_dir = directory.strip().lower()
    if not clean_dir or any(x in clean_dir for x in ("desktop", "downloads", "documents", "user")):
        return os.path.join(UPLOAD_DIR, filename)
    
    # Create subfolders inside UPLOAD_DIR if needed
    subfolder = re.sub(r'[^a-zA-Z0-9_\-]', '', directory)
    target_dir = os.path.join(UPLOAD_DIR, subfolder)
    os.makedirs(target_dir, exist_ok=True)
    return os.path.join(target_dir, filename)

def _inline_fmt(paragraph, text: str):
    """Add text with **bold** and *italic* formatting parsing."""
    parts = re.split(r'\*\*(.+?)\*\*', text)
    for i, part in enumerate(parts):
        if i % 2 == 1:
            paragraph.add_run(part).bold = True
        else:
            for j, ipart in enumerate(re.split(r'\*(.+?)\*', part)):
                if not ipart:
                    continue
                run = paragraph.add_run(ipart)
                if j % 2 == 1:
                    run.italic = True

def _md_to_docx(doc, content: str):
    """Convert markdown text into python-docx paragraphs and tables."""
    lines = content.strip().split("\n")
    i = 0
    while i < len(lines):
        s = lines[i].strip()

        # Markdown table detection
        if "|" in s and i + 1 < len(lines) and re.match(r"^[\|\s\-:]+$", lines[i + 1].strip()):
            table_lines = [s]
            j = i + 2
            while j < len(lines) and "|" in lines[j]:
                table_lines.append(lines[j].strip())
                j += 1
            headers = [c.strip() for c in table_lines[0].split("|") if c.strip()]
            rows = [[c.strip() for c in r.split("|") if c.strip()] for r in table_lines[1:]]
            if headers:
                tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                tbl.style = "Table Grid"
                hdr = tbl.rows[0]
                for ci, h in enumerate(headers):
                    cell = hdr.cells[ci]
                    cell.text = h
                    cell.paragraphs[0].runs[0].bold = True
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row[:len(headers)]):
                        tbl.rows[ri + 1].cells[ci].text = val
            i = j
            continue

        if not s:
            doc.add_paragraph()
        elif s.startswith("#### "):
            doc.add_heading(s[5:], level=4)
        elif s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith("- ") or s.startswith("* "):
            _inline_fmt(doc.add_paragraph(style="List Bullet"), s[2:])
        elif re.match(r"^\d+\.\s", s):
            _inline_fmt(doc.add_paragraph(style="List Number"), re.sub(r"^\d+\.\s", "", s))
        elif s in ("---", "***", "___"):
            doc.add_page_break()
        else:
            _inline_fmt(doc.add_paragraph(), s)
        i += 1

def _make_excel(path: str, content: str):
    wb = openpyxl.Workbook()
    ws = wb.active
    header_done = False
    row_num = 1
    for line in content.strip().split("\n"):
        line = line.strip()
        if not line or re.match(r"^\|?[-:]+\|", line):
            continue
        cells = ([c.strip() for c in line.split("|") if c.strip()]
                 if line.startswith("|") else
                 [c.strip() for c in line.split(",")])
        for col, val in enumerate(cells, 1):
            cell = ws.cell(row=row_num, column=col, value=val)
            if not header_done:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="2E75B6", end_color="2E75B6", fill_type="solid")
                cell.alignment = Alignment(horizontal="center")
        if not header_done:
            header_done = True
        row_num += 1
    for col in ws.columns:
        w = max((len(str(c.value or "")) for c in col), default=8)
        ws.column_dimensions[col[0].column_letter].width = min(w + 4, 50)
    wb.save(path)

def _make_pptx(path: str, content: str):
    prs = Presentation()
    content_layout = prs.slide_layouts[1]
    
    def _new_slide(title_text: str = "", is_title_slide: bool = False):
        layout = prs.slide_layouts[0] if is_title_slide else content_layout
        slide = prs.slides.add_slide(layout)
        if title_text:
            try:
                slide.shapes.title.text = title_text
            except Exception:
                pass
        return slide

    lines = content.strip().split("\n")
    current_slide = None
    current_tf = None
    first_slide = True

    for line in lines:
        s = line.strip()

        if s in ("---", "***"):
            current_slide = None
            current_tf = None
            continue

        if s.startswith("# "):
            title_text = s[2:].strip()
            current_slide = _new_slide(title_text, is_title_slide=first_slide)
            first_slide = False
            current_tf = None
            for ph in current_slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    current_tf = ph.text_frame
                    current_tf.clear()
                    break
            continue

        if s.startswith("## "):
            sub_text = s[3:].strip()
            if current_slide is None:
                current_slide = _new_slide(sub_text)
                first_slide = False
            else:
                for ph in current_slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        tf = ph.text_frame
                        p = tf.add_paragraph()
                        p.text = sub_text
                        p.font.bold = True
                        break
            continue

        if s.startswith("- ") or s.startswith("* "):
            bullet_text = s[2:]
            if current_slide is None:
                current_slide = _new_slide("Slide")
                first_slide = False
            if current_tf is None:
                for ph in current_slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        current_tf = ph.text_frame
                        current_tf.clear()
                        break
            if current_tf:
                p = current_tf.add_paragraph()
                p.text = bullet_text
                p.level = 0
            continue

        if s and current_slide is not None and current_tf is not None:
            p = current_tf.add_paragraph()
            p.text = s

    if not prs.slides._sldIdLst:
        _new_slide("Presentation")
    prs.save(path)

# --- Server-Side create_file ---
def server_create_file(filename: str, directory: str, content: str, file_type: str = "auto") -> str:
    """
    Creates a file on the server. If successful, returns a JSON string containing the
    success status and the web URL for downloading the file.
    """
    # Infer type
    ext = os.path.splitext(filename.lower())[1]
    ftype = file_type.lower()
    
    if ftype == "auto":
        if ext == ".docx": ftype = "word"
        elif ext == ".xlsx": ftype = "excel"
        elif ext == ".pptx": ftype = "pptx"
        elif ext == ".csv": ftype = "csv"
        elif ext == ".pdf": ftype = "pdf"
        elif ext in (".txt", ".md", ".json", ".html", ".js", ".py", ".ts", ".css"): ftype = "text"
        else: ftype = "text"

    # Make sure extension is correct
    if ext == "":
        if ftype == "word": filename += ".docx"
        elif ftype == "excel": filename += ".xlsx"
        elif ftype == "pptx": filename += ".pptx"
        elif ftype == "csv": filename += ".csv"
        elif ftype == "pdf": filename += ".pdf"
        else: filename += ".txt"

    full_path = _resolve_server_path(filename, directory)
    
    try:
        if ftype == "word":
            if not HAS_DOCX:
                return "Error: python-docx library is not installed on the server."
            doc = Document()
            _md_to_docx(doc, content)
            doc.save(full_path)
        elif ftype == "excel":
            if not HAS_OPENPYXL:
                return "Error: openpyxl library is not installed on the server."
            _make_excel(full_path, content)
        elif ftype == "pptx":
            if not HAS_PPTX:
                return "Error: python-pptx library is not installed on the server."
            _make_pptx(full_path, content)
        elif ftype == "csv" or ftype == "text" or ftype == "pdf":
            # For PDF, if it is raw text, we write as text (for PDF conversion or direct read)
            with open(full_path, "w", encoding="utf-8") as f:
                f.write(content)
        else:
            with open(full_path, "w", encoding="utf-8") as f:
                f.write(content)

        # Build download URL (relative to root domain /static/)
        relative_url = f"/static/{os.path.basename(full_path)}"
        # Note: the full URL can be built by the client, or we return the relative path.
        return json.dumps({
            "status": "success",
            "message": f"Created '{filename}' successfully.",
            "file_path": full_path,
            "download_url": relative_url
        })
    except Exception as e:
        return f"Error creating file '{filename}': {str(e)}"

# --- Server-Side convert_file ---
def server_convert_file(file_path: str, target_format: str, output_directory: str = "") -> str:
    """
    Converts a file to a different format on the server.
    Supported:
      - .xlsx -> .csv
      - .csv -> .xlsx
      - .pdf -> .txt
      - .docx -> .txt (extractions)
    """
    # Resolve source path
    src_filename = os.path.basename(file_path)
    if not os.path.exists(file_path):
        # Check if it lives in UPLOAD_DIR
        candidate = os.path.join(UPLOAD_DIR, src_filename)
        if os.path.exists(candidate):
            file_path = candidate
        else:
            return f"Error: Source file '{file_path}' not found on server."

    src_ext = os.path.splitext(file_path.lower())[1]
    tgt_ext = target_format.lower().strip(".")
    if tgt_ext == "word": tgt_ext = "docx"
    elif tgt_ext == "excel": tgt_ext = "xlsx"

    out_name = os.path.splitext(src_filename)[0] + "." + tgt_ext
    out_path = _resolve_server_path(out_name, output_directory)

    try:
        if src_ext == ".xlsx" and tgt_ext == "csv":
            if not HAS_OPENPYXL: return "Error: openpyxl is not installed."
            wb = openpyxl.load_workbook(file_path)
            ws = wb.active
            import csv
            with open(out_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow([v if v is not None else "" for v in row])
            
        elif src_ext == ".csv" and tgt_ext == "xlsx":
            if not HAS_OPENPYXL: return "Error: openpyxl is not installed."
            import csv
            wb = openpyxl.Workbook()
            ws = wb.active
            with open(file_path, encoding="utf-8") as f:
                for r, row in enumerate(csv.reader(f), 1):
                    for c, val in enumerate(row, 1):
                        ws.cell(row=r, column=c, value=val)
            wb.save(out_path)
            
        elif src_ext == ".pdf" and tgt_ext == "txt":
            if not HAS_PYPDF: return "Error: pypdf is not installed."
            reader = pypdf.PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(text)
                
        elif src_ext == ".docx" and tgt_ext == "txt":
            if not HAS_DOCX: return "Error: python-docx is not installed."
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(text)
        else:
            return f"Error: Conversion from {src_ext} to .{tgt_ext} is not supported server-side without office software."

        relative_url = f"/static/{os.path.basename(out_path)}"
        return json.dumps({
            "status": "success",
            "message": f"Converted '{src_filename}' to '{out_name}' successfully.",
            "file_path": out_path,
            "download_url": relative_url
        })
    except Exception as e:
        return f"Error converting file: {str(e)}"
